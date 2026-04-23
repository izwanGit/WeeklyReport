# coding: utf-8
import streamlit as st
import io
import traceback
import sys
import os
import base64
import datetime
import re

try:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

if getattr(sys, 'frozen', False):
    BASE_DIR = sys._MEIPASS
else:
    BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

# -- Page Config --
st.set_page_config(
    page_title="Monthly Report | PETRONAS",
    page_icon=os.path.join(BASE_DIR, "PETRONAS_LOGO_SQUARE.png"),
    layout="wide",
)

# -- Branding Helpers --
def _image_to_data_uri(path, mime_type):
    try:
        with open(os.path.join(BASE_DIR, path), 'rb') as f:
            data = f.read()
        return f"data:{mime_type};base64,{base64.b64encode(data).decode()}"
    except:
        return ""

_logo_square_uri = _image_to_data_uri("PETRONAS_LOGO_SQUARE.png", "image/png")
_logo_sidebar_uri = _image_to_data_uri("PETRONAS_LOGO_HORIZONTAL.svg", "image/svg+xml")


# -- Premium Corporate CSS --
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');

    html, body, [data-testid="stAppViewContainer"] {
        font-family: 'Inter', sans-serif !important;
    }

    [data-testid="stDecoration"],
    [data-testid="stStatusWidget"],
    [data-testid="stSidebarNav"],
    .stDeployButton {
        display: none !important;
        visibility: hidden !important;
    }
    header[data-testid="stHeader"] {
        background: transparent !important;
        border-bottom: none !important;
    }

    [data-testid="stAppViewContainer"] > .main { transition: none !important; }

    [data-testid="stSidebar"] { animation: none !important; }
    [data-testid="stSidebarNav"],
    [data-testid="stSidebarNavItems"],
    [data-testid="stSidebarNavSeparator"],
    [data-testid="stStatusWidget"] {
        display: none !important;
        visibility: hidden !important;
    }
    header[data-testid="stHeader"] { background: #F8FAFC !important; }
    [data-testid="stSidebar"] { border-right: none !important; }

    .main .block-container {
        padding-top: 1rem !important;
        max-width: 1200px !important;
    }
    [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2, [data-testid="stSidebar"] h3 {
        color: #00B1A9 !important; font-weight: 700 !important;
    }
    .stButton > button, .stDownloadButton > button {
        background: #00B1A9 !important; color: white !important;
        border: none !important; border-radius: 10px !important;
        font-weight: 600 !important; transition: all 0.3s ease !important;
        padding: 0.6rem 1.4rem !important;
    }
    .stButton > button:hover, .stDownloadButton > button:hover {
        background: #008C86 !important; transform: translateY(-1px) !important; color: white !important;
    }
    [data-testid="stFileUploader"] {
        border: 2px dashed rgba(0, 177, 169, 0.35) !important;
        border-radius: 12px !important; padding: 16px 20px !important;
    }
    [data-testid="stMetric"] {
        background: #FFFFFF !important; border: 1px solid #E2E8F0 !important;
        border-left: 4px solid #00B1A9 !important; border-radius: 12px !important;
        padding: 1rem 1.2rem !important; box-shadow: 0 2px 8px rgba(0,0,0,0.04) !important;
    }
    [data-testid="stMetricValue"] { color: #00B1A9 !important; font-weight: 800 !important; }
    [data-testid="stMetricLabel"] { color: #4A5568 !important; font-weight: 500 !important; }

    [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] p,
    [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] h1,
    [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] h2,
    [data-testid="stSidebar"] [data-testid="stMarkdownContainer"] h3 {
        margin-bottom: 0px !important;
    }
    [data-testid="stSidebar"] .stTextInput { margin-bottom: -10px !important; }

    .stDeployButton, [data-testid="stDeployButton"], [data-testid="stAppDeployButton"] { display: none !important; }
    #MainMenu { visibility: hidden; } footer { visibility: hidden; }

    .section-label {
        font-size: 0.7rem;
        font-weight: 700;
        text-transform: uppercase;
        letter-spacing: 1px;
        color: #94A3B8;
        margin-bottom: 12px;
        margin-top: 8px;
    }

    [data-testid="stSidebarNav"] { display: none !important; }

    .sidebar-nav {
        display: block;
        padding: 10px 16px;
        margin: 3px 12px;
        border-radius: 8px;
        text-decoration: none !important;
        color: #334155 !important;
        font-family: 'Inter', sans-serif;
        font-weight: 600;
        font-size: 0.88rem;
        transition: all 0.2s ease;
    }
    .sidebar-nav:hover {
        background: #F1F5F9;
        color: #1E293B !important;
        text-decoration: none !important;
    }
    .sidebar-nav.active {
        background: linear-gradient(135deg, #00B1A9, #008C86);
        color: white !important;
        font-weight: 700;
        box-shadow: 0 4px 12px rgba(0,177,169,0.25);
    }
    .sidebar-sep { border: none; border-top: 1px solid #E2E8F0; margin: 16px 12px; }
    .genie-link {
        font-size: 0.85rem; font-weight: 500;
        color: #31333F !important; text-decoration: none !important;
        transition: all 0.2s ease !important; cursor: pointer !important;
    }
    .genie-link:hover { color: #00B1A9 !important; text-decoration: none !important; }
</style>
""", unsafe_allow_html=True)


# -- Sidebar --
with st.sidebar:
    st.markdown(f"""
<div style="text-align:center; padding:8px 0 20px 0;">
    <a href="/" target="_self" style="display:inline-block;">
        <img src="{_logo_sidebar_uri}" style="height:56px; transition: transform 0.2s; cursor: pointer;" onmouseover="this.style.transform='scale(1.05)'" onmouseout="this.style.transform='scale(1)'"/>
    </a>
</div>
""", unsafe_allow_html=True)

    st.markdown("### Report Settings")
    _months = ["January", "February", "March", "April", "May", "June",
               "July", "August", "September", "October", "November", "December"]
    _now = datetime.date.today()
    sel_month = st.selectbox("Report Month", options=_months, index=_now.month - 1)
    _current_year = _now.year
    _years = [str(y) for y in range(_current_year - 1, _current_year + 3)]
    sel_year = st.selectbox("Report Year", options=_years, index=1)

    st.markdown("<div style='margin-top: -10px;'></div>", unsafe_allow_html=True)
    st.markdown("### Data Upload")
    st.markdown("<a href='https://app.powerbi.com/groups/81a248dd-b149-45b3-9af2-2f0206f1df7b/reports/6b253891-d1d6-4570-b6fb-ef0e57214ee5/f5dbf7be18699d9dedd8?experience=power-bi' target='_blank' class='genie-link'>Power BI PDF Export ↗</a>", unsafe_allow_html=True)
    pdf_file = st.file_uploader("Power BI PDF Export", type=['pdf'], label_visibility="collapsed")

    TEMPLATE_PATH = os.path.join(BASE_DIR, "template.pptx")
    pptx_available_locally = os.path.exists(TEMPLATE_PATH)
    uploaded_template = None
    if not pptx_available_locally:
        st.warning("template.pptx not found in app folder.")
        uploaded_template = st.file_uploader("Upload PPTX Template", type=['pptx'])
    else:
        st.success("template.pptx detected.")

st.markdown("""
<a href="/" target="_self" style="text-decoration: none; display: inline-flex; align-items: center; gap: 8px; font-weight: 600; color: #64748B; margin-bottom: 16px; transition: color 0.2s ease;">
    <svg width="18" height="18" viewBox="0 0 24 24" fill="none" stroke="currentColor" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><line x1="19" y1="12" x2="5" y2="12"></line><polyline points="12 19 5 12 12 5"></polyline></svg>
    Back to Hub
</a>
""", unsafe_allow_html=True)

# -- Header Banner --
st.markdown(f"""
<style>
.banner-title {{ color: #FFFFFF !important; text-transform: uppercase !important; font-weight: 800 !important; text-shadow: 0px 2px 4px rgba(0,0,0,0.3) !important; margin: 0 !important; line-height: 1.1 !important; white-space: nowrap; font-size: clamp(1.2rem, 3.5vw, 1.8rem) !important; letter-spacing: 0.1px; }}
.banner-subtitle {{ color: #FFFFFF !important; font-weight: 400 !important; text-shadow: 0px 1px 3px rgba(0,0,0,0.2) !important; margin: 4px 0 0 0 !important; white-space: nowrap; font-size: clamp(0.85rem, 2vw, 1.0rem) !important; opacity: 0.95 !important; }}
</style>
<div style="display: flex; align-items: center; gap: 24px; padding: 22px 32px; background-color: #00B1A9; border-radius: 20px; margin-bottom: 2rem; box-shadow: 0 12px 35px rgba(0, 177, 169, 0.25); overflow: hidden; border: 1px solid rgba(255, 255, 255, 0.15);">
<img src="{_logo_square_uri}" style="height: 80px; flex-shrink: 0; filter: drop-shadow(1px 1px 0 white) drop-shadow(-1px -1px 0 white) drop-shadow(1px -1px 0 white) drop-shadow(-1px 1px 0 white);" />
<div style="min-width: 0;">
<h1 class="banner-title">Monthly PPTX Automation</h1>
<p class="banner-subtitle">Power BI dashboard export to corporate PowerPoint deck - zero-touch pipeline.</p>
</div>
</div>
""", unsafe_allow_html=True)

# -- Dependency Check --
if not PPTX_AVAILABLE:
    st.error("Required libraries (python-pptx, PyMuPDF) are not installed. Please run: pip install -r requirements.txt")
    st.stop()


# ==============================================================
# HELPER FUNCTIONS
# ==============================================================

def _emu_to_inches(emu):
    """PowerPoint uses English Metric Units: 1 inch = 914400 EMU."""
    return emu / 914400.0


def _auto_crop_bottom(png_bytes, margin=20, threshold=245):
    """Trim blank whitespace from the bottom of a PNG image."""
    from PIL import Image
    img = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    pixels = img.load()
    w, h = img.size
    last_content_row = h - 1
    for y in range(h - 1, -1, -1):
        row_is_blank = True
        for x in range(0, w, 5):
            r, g, b = pixels[x, y]
            if r < threshold or g < threshold or b < threshold:
                row_is_blank = False
                break
        if not row_is_blank:
            last_content_row = y
            break
    crop_bottom = min(last_content_row + margin, h)
    if crop_bottom < h - 30:
        img = img.crop((0, 0, w, crop_bottom))
    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def _get_image_size(png_bytes):
    """Return (width_px, height_px) of a PNG image."""
    from PIL import Image
    with Image.open(io.BytesIO(png_bytes)) as img:
        return img.size


def _calc_fit_rect(img_w, img_h, box_left, box_top, box_w, box_h):
    """
    Fit image inside bounding box preserving aspect ratio, centred.
    Used for CHART replacements (slides 4, 8).
    All box_* values are in EMU; img_w/img_h are pixels (ratio only matters).
    Returns (left, top, width, height) all in EMU.
    """
    scale = min(box_w / img_w, box_h / img_h)
    new_w = int(img_w * scale)
    new_h = int(img_h * scale)
    offset_x = (box_w - new_w) // 2
    offset_y = (box_h - new_h) // 2
    return box_left + offset_x, box_top + offset_y, new_w, new_h


def _calc_list_rect(img_w, img_h, box_left, box_top, box_w):
    """
    Width-constrained placement for variable-length module lists (slides 6, 10).
    Locks to the bounding box width and lets height follow the actual
    (already auto-cropped) image content.
    Used so that a month with 5 modules is short and one with 20 is tall,
    with no stretching either way.
    Returns (left, top, width, height) all in EMU.
    """
    scale = box_w / img_w
    new_h = int(img_h * scale)
    return box_left, box_top, box_w, new_h


def _extract_month_values_spatial(page):
    """
    Extract (month_label, numeric_value) pairs from chart using spatial analysis.

    IMPORTANT RULE:
    - Values are returned in LEFT -> RIGHT order exactly as they appear on chart.
    - This allows summary logic to always use:
        values[-1] = latest (rightmost)
        values[-2] = previous (immediately beside latest)

    Handles labels like:
      - Feb 2026
      - Feb 26
      - Feb '26
      - Feb 26'
    """
    DATE_PAT = re.compile(
        r'^(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)'
        r'\s+(?:\d{4}|\'?\d{2}\'?)$',
        re.IGNORECASE
    )

    try:
        blocks = page.get_text("dict")["blocks"]
    except Exception:
        return []

    spans = []
    for block in blocks:
        if "lines" not in block:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                text = span["text"].strip()
                if not text:
                    continue
                x0, y0, x1, y1 = span["bbox"]
                spans.append({
                    "text": text,
                    "x0": x0,
                    "x1": x1,
                    "y0": y0,
                    "y1": y1,
                    "cx": (x0 + x1) / 2,
                    "cy": (y0 + y1) / 2
                })

    if not spans:
        return []

    # ----------------------------------------------------------
    # Step 1: find month/year labels on x-axis
    # ----------------------------------------------------------
    month_spans = []
    for sp in spans:
        txt = _clean_cell_text(sp["text"])
        if DATE_PAT.match(txt):
            month_spans.append({
                "label": txt,
                "cx": sp["cx"],
                "cy": sp["cy"],
                "x0": sp["x0"],
                "x1": sp["x1"]
            })

    if not month_spans:
        return []

    # Group labels by y-row; choose the lowest row with the most month labels
    rows = []
    Y_TOL = 14

    for ms in sorted(month_spans, key=lambda s: s["cy"]):
        placed = False
        for grp in rows:
            if abs(ms["cy"] - grp["y"]) <= Y_TOL:
                grp["items"].append(ms)
                grp["y"] = sum(i["cy"] for i in grp["items"]) / len(grp["items"])
                placed = True
                break
        if not placed:
            rows.append({"y": ms["cy"], "items": [ms]})

    # Best row = most labels, then lower on page (larger y)
    rows.sort(key=lambda g: (len(g["items"]), g["y"]), reverse=True)
    axis_months = sorted(rows[0]["items"], key=lambda s: s["cx"])

    if len(axis_months) < 2:
        return []

    # ----------------------------------------------------------
    # Step 2: extract number nearest above each month label
    # ----------------------------------------------------------
    # Estimate x tolerance from spacing between months
    gaps = []
    for i in range(1, len(axis_months)):
        gaps.append(axis_months[i]["cx"] - axis_months[i - 1]["cx"])

    if gaps:
        avg_gap = sum(gaps) / len(gaps)
        x_tol = max(28, min(90, avg_gap * 0.45))
    else:
        x_tol = 40

    results = []

    for m in axis_months:
        best_val = None
        best_score = None

        for sp in spans:
            txt = _clean_cell_text(sp["text"])

            # skip month labels themselves
            if DATE_PAT.match(txt):
                continue

            # only numbers
            clean = txt.replace(",", "").replace("%", "").strip()
            nm = re.fullmatch(r'-?(\d+(?:\.\d+)?)', clean)
            if not nm:
                continue

            val = float(nm.group(1))

            # must be above the month label
            if sp["cy"] >= m["cy"] - 2:
                continue

            x_dist = abs(sp["cx"] - m["cx"])
            if x_dist > x_tol:
                continue

            # prefer number closest vertically above the month label
            vertical_dist = m["cy"] - sp["cy"]

            # soft filter: ignore values too far up (often axis ticks)
            if vertical_dist > 220:
                continue

            score = (vertical_dist, x_dist)

            if best_score is None or score < best_score:
                best_score = score
                best_val = val

        if best_val is not None:
            results.append((m["label"], best_val))

    # IMPORTANT: return exactly in left-to-right chart order
    return results


def _extract_month_values_text(page_text):
    """
    Fallback text extractor.

    IMPORTANT:
    - Returns values in order of appearance from the PDF text,
      not Jan..Dec hard-coded order.
    - This keeps alignment with chart left -> right order as much as possible.
    """
    DATE_PAT = re.compile(
        r'^(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)'
        r'\s+(?:\d{4}|\'?\d{2}\'?)$',
        re.IGNORECASE
    )

    lines = [_clean_cell_text(l) for l in page_text.split('\n') if _clean_cell_text(l)]
    found = []
    used_labels = set()

    for i, line in enumerate(lines):
        if not DATE_PAT.match(line):
            continue

        label = line
        if label in used_labels:
            continue

        best_val = None

        # Prefer number AFTER label
        for j in range(i + 1, min(i + 6, len(lines))):
            clean = lines[j].replace(",", "").replace("%", "").strip()
            nm = re.fullmatch(r'-?(\d+(?:\.\d+)?)', clean)
            if nm:
                best_val = float(nm.group(1))
                break

        # If not found, try BEFORE label
        if best_val is None:
            for j in range(i - 1, max(i - 6, -1), -1):
                clean = lines[j].replace(",", "").replace("%", "").strip()
                nm = re.fullmatch(r'-?(\d+(?:\.\d+)?)', clean)
                if nm:
                    best_val = float(nm.group(1))
                    break

        if best_val is not None:
            found.append((label, best_val))
            used_labels.add(label)

    return found


def _normalize_chart_order(vals):
    """
    Enforce LEFT → RIGHT chronological order before summary logic.
    Sorts by (year, month_number) so that:
      values[-1] = rightmost / latest month
      values[-2] = month immediately before it
    This fixes the root cause where PDF text extraction returns
    months in arbitrary order (e.g. Apr before Feb).
    """
    if not vals or len(vals) < 2:
        return vals

    MONTH_INDEX = {
        'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4,
        'may': 5, 'jun': 6, 'jul': 7, 'aug': 8,
        'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12
    }

    def _key(item):
        label, _ = item
        parts = label.lower().replace("'", "").split()
        if len(parts) >= 2:
            m = MONTH_INDEX.get(parts[0][:3], 0)
            try:
                y = int(parts[1])
                # Handle 2-digit years (e.g. "26" -> 2026)
                if y < 100:
                    y += 2000
            except ValueError:
                y = 0
            return (y, m)
        return (0, 0)

    return sorted(vals, key=_key)


def _extract_with_fallback(pdf, page_num, key, log):
    """
    Try spatial extraction first, fall back to text-based.
    page_num is 1-indexed.
    Returns list of (month, value) tuples.
    """
    if page_num > len(pdf):
        return []

    page = pdf.load_page(page_num - 1)

    # Strategy 1: Spatial analysis (most accurate)
    vals = _extract_month_values_spatial(page)
    if vals:
        log.append(f"INFO | {key}: Spatial extraction OK ({len(vals)} points)")
        if len(vals) >= 2:
            log.append(
                f"INFO | {key}: latest comparison pair = "
                f"{vals[-2][0]}={int(round(vals[-2][1]))} -> {vals[-1][0]}={int(round(vals[-1][1]))}"
            )
        return vals

    # Strategy 2: Text-based fallback
    raw_text = page.get_text()
    vals = _extract_month_values_text(raw_text)
    if vals:
        log.append(f"INFO | {key}: Text fallback OK ({len(vals)} points)")
        if len(vals) >= 2:
            log.append(
                f"INFO | {key}: latest comparison pair = "
                f"{vals[-2][0]}={int(round(vals[-2][1]))} -> {vals[-1][0]}={int(round(vals[-1][1]))}"
            )
        return vals

    # Nothing worked - dump raw text for debugging
    snippet = raw_text[:300].replace('\n', ' | ')
    log.append(f"INFO | {key}: NO DATA. Raw: '{snippet}...'")
    return []


# ==============================================================
# TABLE EXTRACTION HELPERS (FIXED)
# ==============================================================

PLACEHOLDER_VALUES = {"", "-", "—", "□", "☐", "☑", "■"}

def _clean_cell_text(text):
    """Normalize whitespace and strip filler characters."""
    if text is None:
        return ""
    text = re.sub(r"\s+", " ", str(text)).strip()
    return text

def _is_placeholder_value(text):
    text = _clean_cell_text(text)
    return text in PLACEHOLDER_VALUES

def _is_integer_text(text):
    text = _clean_cell_text(text).replace(",", "")
    return bool(re.fullmatch(r"\d+", text))

def _row_has_real_content(category, root_cause, action_plan, total_tickets):
    """
    Accept only rows that have real business data.
    A valid row should have at least:
      - a meaningful category, OR
      - a meaningful root cause, OR
      - a numeric ticket count
    """
    category = _clean_cell_text(category)
    root_cause = _clean_cell_text(root_cause)
    action_plan = _clean_cell_text(action_plan)
    total_tickets = _clean_cell_text(total_tickets)

    has_category = category and not _is_placeholder_value(category)
    has_root = root_cause and not _is_placeholder_value(root_cause)
    has_total = _is_integer_text(total_tickets)

    return has_category or has_root or has_total


def _extract_root_cause_table(pdf, page_num, label, log):
    """
    Extract Root Cause table from PDF page.

    Expected PPT layout:
    Month | [blank header / Product Categorization Tier 3 values] |
    Root Cause | Action Plans | Total Tickets
    """
    MONTH_PAT = re.compile(
        r'(Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|'
        r'Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)'
        r'\s+\d{4}', re.IGNORECASE
    )

    if page_num > len(pdf):
        log.append(f"WARN | {label} table: page {page_num} not in PDF")
        return []

    page = pdf.load_page(page_num - 1)
    page_w = page.rect.width

    try:
        blocks = page.get_text("dict")["blocks"]
    except Exception as e:
        log.append(f"WARN | {label} table: failed to get text dict ({e})")
        return []

    # Collect all spans with positions
    spans = []
    for block in blocks:
        if "lines" not in block:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                text = span["text"].strip()
                if not text:
                    continue
                x0, y0, x1, y1 = span["bbox"]
                spans.append({
                    "text": text,
                    "x0": x0,
                    "x1": x1,
                    "cx": (x0 + x1) / 2,
                    "cy": (y0 + y1) / 2
                })

    if not spans:
        log.append(f"WARN | {label} table: no spans found")
        return []

    # Group spans into rows by Y
    rows_by_y = []
    sorted_spans = sorted(spans, key=lambda s: s["cy"])
    current_row = []
    current_y = None
    Y_TOL = 8

    for sp in sorted_spans:
        if current_y is None:
            current_row = [sp]
            current_y = sp["cy"]
        elif abs(sp["cy"] - current_y) <= Y_TOL:
            current_row.append(sp)
        else:
            current_row.sort(key=lambda s: s["x0"])
            rows_by_y.append((current_y, current_row))
            current_row = [sp]
            current_y = sp["cy"]

    if current_row:
        current_row.sort(key=lambda s: s["x0"])
        rows_by_y.append((current_y, current_row))

    # Find header row
    header_y = None
    header_row = None
    for cy, row in rows_by_y:
        row_text = " ".join(s["text"] for s in row).lower()
        if "root cause" in row_text and "action" in row_text and "total" in row_text:
            header_y = cy
            header_row = row
            break

    if header_y is None:
        log.append(f"WARN | {label} table: header row not found")
        return []

    # Find month text
    month_text = ""
    for sp in spans:
        m = MONTH_PAT.search(sp["text"])
        if m:
            month_text = m.group(0)
            break

    # ==========================================================
    # Derive 5 body columns:
    # 0 = Month
    # 1 = Product Categorization Tier 3 (blank header in PPT)
    # 2 = Root Cause
    # 3 = Action Plans
    # 4 = Total Tickets
    # ==========================================================
    positions = {}
    for sp in header_row:
        t = sp["text"].lower()
        xf0 = sp["x0"] / page_w
        xf1 = sp["x1"] / page_w

        if "month" in t:
            positions["month_x0"] = xf0
            positions["month_x1"] = xf1
        elif "root" in t or "cause" in t:
            positions["root_x0"] = xf0
            positions["root_x1"] = xf1
        elif "action" in t:
            positions["action_x0"] = xf0
            positions["action_x1"] = xf1
        elif "total" in t or "ticket" in t:
            positions["total_x0"] = xf0
            positions["total_x1"] = xf1

    # Conservative fallback boundaries
    month_end = positions.get("month_x1", 0.12)
    root_start = positions.get("root_x0", 0.26)
    action_start = positions.get("action_x0", 0.62)
    total_start = positions.get("total_x0", 0.88)

    # Ensure sane ordering
    if root_start <= month_end:
        root_start = max(month_end + 0.10, 0.26)
    if action_start <= root_start:
        action_start = max(root_start + 0.20, 0.62)
    if total_start <= action_start:
        total_start = max(action_start + 0.15, 0.88)

    # 5 columns:
    # month | tier3 | root cause | action | total
    col_bounds = [0.00, month_end, root_start, action_start, total_start, 1.01]
    log.append(f"INFO | {label} col bounds: {[round(x, 3) for x in col_bounds]}")

    # Collect candidate rows below header
    candidate_rows = []
    for cy, row in rows_by_y:
        if cy <= header_y + 5:
            continue

        row_texts = [_clean_cell_text(s["text"]) for s in row]
        non_month_texts = [t for t in row_texts if t and not MONTH_PAT.search(t)]

        if not non_month_texts:
            continue  # Skip month-only row

        candidate_rows.append((cy, row))

    if not candidate_rows:
        log.append(f"WARN | {label} table: no candidate rows below header")
        return []

    # Parse each candidate row into columns
    results = []

    for cy, row in candidate_rows:
        cols = ["", "", "", "", ""]

        for sp in row:
            txt = _clean_cell_text(sp["text"])
            if not txt:
                continue
            if MONTH_PAT.search(txt):
                continue

            # Use center X for better placement when text is centered in cells
            x_frac = sp["cx"] / page_w

            for col_idx in range(5):
                if col_bounds[col_idx] <= x_frac < col_bounds[col_idx + 1]:
                    cols[col_idx] += (" " if cols[col_idx] else "") + txt
                    break

        cols = [_clean_cell_text(c) for c in cols]

        # Find rightmost integer as total tickets
        total = ""
        for sp in sorted(row, key=lambda s: s["x0"], reverse=True):
            clean = _clean_cell_text(sp["text"]).replace(",", "")
            if re.fullmatch(r"\d+", clean):
                total = clean
                break

        product_tier_3 = cols[1]
        root_cause = cols[2]
        action_plan = cols[3]
        total_tickets = total if total else cols[4]

        # Remove placeholders
        if _is_placeholder_value(product_tier_3):
            product_tier_3 = ""
        if _is_placeholder_value(root_cause):
            root_cause = ""
        if _is_placeholder_value(action_plan):
            action_plan = ""
        if _is_placeholder_value(total_tickets):
            total_tickets = ""

        # Fallback: derive Product Categorization Tier 3 from Root Cause if needed
        # Example: "Issue Related to Core HR" -> "Core HR"
        if not product_tier_3 and root_cause:
            m = re.search(r'Issue\s+Related\s+to\s+(.+)$', root_cause, re.IGNORECASE)
            if m:
                product_tier_3 = m.group(1).strip().lstrip("-").strip()

        # Optional: normalize root cause to start with "- "
        if root_cause and not root_cause.startswith("-"):
            root_cause = f"- {root_cause}"

        if not _row_has_real_content(product_tier_3, root_cause, action_plan, total_tickets):
            log.append(
                f"INFO | {label} skipped phantom row at y={cy:.1f}: "
                f"tier3='{product_tier_3}' rc='{root_cause}' ap='{action_plan}' total='{total_tickets}'"
            )
            continue

        if not action_plan:
            action_plan = "To identify action plan to reduce tickets"

        row_result = {
            "month": month_text,
            "product_tier_3": product_tier_3,
            "root_cause": root_cause,
            "action_plan": action_plan,
            "total_tickets": total_tickets,
            "_y": cy
        }

        results.append(row_result)

        log.append(
            f"INFO | {label} row kept: "
            f"tier3='{product_tier_3}' | rc='{root_cause}' | "
            f"ap='{action_plan}' | tickets='{total_tickets}'"
        )

    # Keep only top 3 real rows
    results.sort(key=lambda r: r["_y"])
    results = results[:3]

    for r in results:
        r.pop("_y", None)

    if not results:
        log.append(f"WARN | {label} table: no valid rows extracted after filtering")
    else:
        log.append(f"INFO | {label} table: extracted {len(results)} valid rows")

    return results


def _fill_pptx_table(pptx_shape, data_rows, label, sel_month, sel_year, log):
    """
    Fill existing PPTX table with extracted data. Preserves formatting.
    IMPORTANT:
    - Column 0 = Month
    - Column 1 = Product Categorization Tier 3 values (header intentionally left blank)
    - Column 2 = Root Cause
    - Column 3 = Action Plans
    - Column 4 = Total Tickets
    """
    def _set_cell(cell, text):
        tf = cell.text_frame

        if not tf.paragraphs:
            tf.add_paragraph()

        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ""

        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = str(text)
        else:
            run = tf.paragraphs[0].add_run()
            run.text = str(text)

    tbl = pptx_shape.table
    num_rows = len(tbl.rows)
    num_cols = len(tbl.columns)
    log.append(f"INFO | {label} PPTX table: {num_rows} rows x {num_cols} cols")

    month_label = f"{sel_month} {sel_year}"

    # Clear body rows only (leave header row untouched)
    for r in range(1, min(num_rows, 4)):
        for c in range(num_cols):
            try:
                _set_cell(tbl.cell(r, c), "")
            except Exception as e:
                log.append(f"WARN | {label} clear cell ({r},{c}) failed: {e}")

    # Write extracted rows
    for i, row_data in enumerate(data_rows[:3]):
        pptx_row = i + 1
        if pptx_row >= num_rows:
            log.append(f"WARN | {label} table: insufficient PPT rows")
            continue

        try:
            # Month only on first row (merged-look effect)
            if i == 0 and num_cols > 0:
                _set_cell(tbl.cell(pptx_row, 0), month_label)

            # Column 1 body = Product Categorization Tier 3
            # Header remains whatever is already in template (blank)
            if num_cols > 1:
                _set_cell(tbl.cell(pptx_row, 1), row_data.get("product_tier_3", ""))

            if num_cols > 2:
                _set_cell(tbl.cell(pptx_row, 2), row_data.get("root_cause", ""))

            if num_cols > 3:
                _set_cell(tbl.cell(pptx_row, 3), row_data.get("action_plan", ""))

            if num_cols > 4:
                _set_cell(tbl.cell(pptx_row, 4), row_data.get("total_tickets", ""))

            log.append(
                f"  TABLE | {label} row {pptx_row}: "
                f"tier3='{row_data.get('product_tier_3','')}' | "
                f"root='{row_data.get('root_cause','')}' | "
                f"action='{row_data.get('action_plan','')}' | "
                f"total='{row_data.get('total_tickets','')}'"
            )

        except Exception as e:
            log.append(f"WARN | {label} table row {pptx_row} error: {e}")


# ==============================================================
# DATE AND SUMMARY UPDATE HELPERS (NEW & PATCHED)
# ==============================================================

def _shape_text(shape):
    """Return full text from a shape text frame."""
    if not hasattr(shape, "text_frame"):
        return ""
    parts = []
    for p in shape.text_frame.paragraphs:
        txt = "".join(r.text for r in p.runs).strip()
        if txt:
            parts.append(txt)
    return "\n".join(parts).strip()


def _set_paragraph_text(para, text):
    """Replace paragraph text while preserving paragraph formatting/bullets."""
    if not para.runs:
        para.add_run()
    para.runs[0].text = str(text)
    for r in para.runs[1:]:
        r.text = ""


def _find_summary_title_shape(slide):
    """
    Find the text shape that is the 'Summary' heading.
    Prefers exact match, then short text containing summary.
    """
    candidates = []

    for sh in slide.shapes:
        if not hasattr(sh, "text_frame"):
            continue

        txt = _shape_text(sh).strip().lower()
        if not txt:
            continue

        if txt == "summary":
            candidates.append((0, sh.top, sh.left, sh))
        elif "summary" in txt and len(txt) <= 30:
            candidates.append((1, sh.top, sh.left, sh))

    if not candidates:
        return None

    candidates.sort(key=lambda x: (x[0], x[1], x[2]))
    return candidates[0][3]


def _find_summary_body_shape(slide, summary_shape):
    """
    Find the textbox most likely to contain the summary bullets.
    Assumption: it is below the Summary heading and roughly left-aligned with it.
    """
    if summary_shape is None:
        return None

    candidates = []
    MAX_LEFT_DIFF = 3 * 914400   # 3 inches in EMU
    MAX_TOP_GAP   = 4 * 914400   # 4 inches in EMU

    for sh in slide.shapes:
        if sh == summary_shape:
            continue
        if not hasattr(sh, "text_frame"):
            continue

        txt = _shape_text(sh).strip()
        if not txt:
            continue

        # Must be below the summary title
        if sh.top <= summary_shape.top:
            continue

        top_gap = sh.top - summary_shape.top
        left_diff = abs(sh.left - summary_shape.left)

        if top_gap > MAX_TOP_GAP:
            continue

        # Prefer textboxes roughly aligned under Summary
        candidates.append((top_gap, left_diff, sh.top, sh.left, sh))

    if not candidates:
        return None

    # Nearest below, then closest horizontal alignment
    candidates.sort(key=lambda x: (x[0], x[1], x[2], x[3]))
    return candidates[0][4]


def _get_meaningful_paragraphs(tf):
    """
    Return non-empty paragraphs excluding the title 'Summary'.
    """
    paras = []
    for p in tf.paragraphs:
        txt = "".join(r.text for r in p.runs).strip()
        if not txt:
            continue
        if txt.lower() == "summary":
            continue
        paras.append(p)
    return paras


def _parse_month_label(label):
    MONTH_INDEX = {
        'jan': 1, 'feb': 2, 'mar': 3, 'apr': 4,
        'may': 5, 'jun': 6, 'jul': 7, 'aug': 8,
        'sep': 9, 'oct': 10, 'nov': 11, 'dec': 12
    }
    parts = label.lower().replace("'", "").split()
    if len(parts) < 2:
        return None

    mon = MONTH_INDEX.get(parts[0][:3])
    if not mon:
        return None

    try:
        year = int(parts[1])
        if year < 100:
            year += 2000
    except ValueError:
        return None

    return year, mon


def _month_diff(label_a, label_b):
    """
    Returns number of months from a -> b
    Example: Mar 2026 -> Apr 2026 = 1
    """
    a = _parse_month_label(label_a)
    b = _parse_month_label(label_b)
    if not a or not b:
        return None
    return (b[0] - a[0]) * 12 + (b[1] - a[1])


def _replace_in_para(para, pattern, replacement, flags=re.IGNORECASE):
    """Replace text matching pattern in a paragraph, keeping first run's formatting."""
    if not para.runs:
        return False
    full_text = "".join(run.text for run in para.runs)
    if not re.search(pattern, full_text, flags):
        return False
    new_text = re.sub(pattern, replacement, full_text, flags=flags)
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""
    return True


def update_dates_in_pptx(prs, new_month, new_year):
    """
    Replace month/year date references across ALL slides.
    PATCHED: Option B - Safely avoids touching summary body textboxes so 
    template placeholders are preserved for the summary engine.
    """
    MONTH_NAMES = ['January', 'February', 'March', 'April', 'May', 'June',
                   'July', 'August', 'September', 'October', 'November', 'December']
    MONTH_ABBREVS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun',
                     'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    all_months = sorted(MONTH_NAMES + MONTH_ABBREVS, key=len, reverse=True)
    month_pat = '|'.join(re.escape(m) for m in all_months)
    changes = []
    
    for slide_num, slide in enumerate(prs.slides):
        # Identify the summary body on this slide to skip it
        summary_title = _find_summary_title_shape(slide)
        summary_body = _find_summary_body_shape(slide, summary_title) if summary_title else None

        for shape in slide.shapes:
            if summary_body and shape == summary_body:
                continue # Skip summary bullets

            if not hasattr(shape, "text_frame"):
                continue
            for para in shape.text_frame.paragraphs:
                original = "".join(run.text for run in para.runs)
                if _replace_in_para(para, rf'({month_pat})\s+(20\d{{2}})', f'{new_month} {new_year}'):
                    changes.append(f"  DATE | Slide {slide_num + 1}: '{original.strip()[:60]}'")
                    continue
                if _replace_in_para(para, r'\b(20\d{2})\b', str(new_year)):
                    changes.append(f"  DATE | Slide {slide_num + 1}: year updated")
    return changes


def _analyze_trend(values):
    """
    STRICT rule:
    Compare ONLY last two points IF they are adjacent months.
    Returns dict with: direction, diff, pct_change, streak, trend_desc, is_peak, is_trough
    """
    if not values or len(values) < 2:
        return None

    prev_label, prev_val = values[-2]
    curr_label, curr_val = values[-1]

    # Must be adjacent months (e.g. Mar -> Apr)
    gap = _month_diff(prev_label, curr_label)
    if gap != 1:
        return None

    curr = int(round(curr_val))
    prev = int(round(prev_val))
    diff = abs(curr - prev)
    direction = "increased" if curr > prev else "decreased" if curr < prev else "unchanged"

    # Percentage change
    if prev > 0:
        pct = round((diff / prev) * 100)
    else:
        pct = 0

    # Consecutive streak analysis (how many months in same direction)
    streak = 1
    if len(values) >= 3:
        for k in range(len(values) - 2, 0, -1):
            v_curr = values[k][1]
            v_prev = values[k - 1][1]
            if direction == "increased" and v_curr > v_prev:
                streak += 1
            elif direction == "decreased" and v_curr < v_prev:
                streak += 1
            else:
                break

    # Peak / trough detection
    all_vals = [v[1] for v in values]
    is_peak = curr == max(all_vals) and curr > prev
    is_trough = curr == min(all_vals) and curr < prev

    return {
        "curr": curr,
        "prev": prev,
        "diff": diff,
        "pct": pct,
        "direction": direction,
        "streak": streak,
        "is_peak": is_peak,
        "is_trough": is_trough,
    }


def _make_smart_ticket_bullet(vals, sel_month, sel_year):
    """Generate an intelligent ticket trend summary bullet."""
    if not vals or len(vals) < 2:
        return None

    a = _analyze_trend(vals)
    if not a:
        return None

    curr, prev, diff, pct = a["curr"], a["prev"], a["diff"], a["pct"]
    direction = a["direction"]

    if direction == "unchanged":
        return f"Ticket volume remained steady at {curr} in {sel_month} {sel_year}"

    verb = "increased" if direction == "increased" else "decreased"

    # Build the core sentence
    parts = [f"Ticket logged {verb} by {diff}"]

    # Add percentage if meaningful
    if pct > 0 and prev > 0:
        parts[0] += f" ({pct}%)"

    parts[0] += f" in {sel_month} {sel_year}, from {prev} to {curr}"

    # Add streak context if notable
    if a["streak"] >= 3:
        parts.append(f"This marks {a['streak']} consecutive months of {verb[:-1] if verb.endswith('d') else verb}ing trend")
    elif a["streak"] == 2:
        parts.append(f"continuing the {verb[:-1] if verb.endswith('d') else verb}ing trend from last month")

    # Add peak/trough context
    if a["is_peak"]:
        parts.append("reaching the highest level in the observed period")
    elif a["is_trough"]:
        parts.append("reaching the lowest level in the observed period")

    return ". ".join(parts)


def _make_smart_ageing_bullet(vals, sel_month, sel_year):
    """Generate an intelligent ageing trend summary bullet."""
    if not vals or len(vals) < 2:
        return None

    a = _analyze_trend(vals)
    if not a:
        return None

    curr, prev, diff, pct = a["curr"], a["prev"], a["diff"], a["pct"]
    direction = a["direction"]

    if direction == "unchanged":
        if curr == 0:
            return f"No ageing tickets recorded in {sel_month} {sel_year}, maintaining zero backlog"
        return f"Ageing ticket count remained at {curr} in {sel_month} {sel_year}"

    if curr == 0 and prev > 0:
        return f"All ageing tickets resolved - count dropped from {prev} to 0 in {sel_month} {sel_year}"

    verb = "increased" if direction == "increased" else "decreased"
    sign = "+" if direction == "increased" else "-"

    parts = [f"Ageing ticket {verb} from {prev} to {curr} in {sel_month} {sel_year} ({sign}{diff})"]

    # Add percentage context for significant changes
    if pct >= 50 and diff >= 3:
        parts.append(f"a significant {pct}% {'rise' if direction == 'increased' else 'reduction'}")
    elif pct > 0 and prev > 0:
        parts[0] = parts[0].replace(f"({sign}{diff})", f"({sign}{diff}, {pct}%)")

    # Streak context
    if a["streak"] >= 3:
        parts.append(f"continuing a {a['streak']}-month {'upward' if direction == 'increased' else 'downward'} trend")

    return ". ".join(parts)


def update_summary_bullets(prs, sel_month, sel_year,
                           sr_trend_vals, sr_ageing_vals,
                           inc_trend_vals, inc_ageing_vals):
    """
    Update summary bullets on Slides 4 and 8.
    FIXED:
    - Finds the 'Summary' title textbox
    - Finds the actual body textbox below it
    - Replaces the first 2 bullet paragraphs directly
    """
    changes = []

    CONFIG = [
        (3, sr_trend_vals,  sr_ageing_vals,  "SR"),
        (7, inc_trend_vals, inc_ageing_vals, "INC"),
    ]

    for slide_idx, ticket_vals, ageing_vals, label in CONFIG:
        if slide_idx >= len(prs.slides):
            changes.append(f"  SUMM | Slide {slide_idx+1}: slide missing")
            continue

        slide = prs.slides[slide_idx]

        ticket_bullet = _make_smart_ticket_bullet(ticket_vals, sel_month, sel_year)
        ageing_bullet = _make_smart_ageing_bullet(ageing_vals, sel_month, sel_year)

        summary_title = _find_summary_title_shape(slide)
        if not summary_title:
            changes.append(f"  SUMM | Slide {slide_idx+1}: Summary title not found")
            continue

        summary_body = _find_summary_body_shape(slide, summary_title)
        if not summary_body:
            changes.append(f"  SUMM | Slide {slide_idx+1}: Summary body textbox not found")
            continue

        tf = summary_body.text_frame
        paras = _get_meaningful_paragraphs(tf)

        # If template has fewer than 2 meaningful paragraphs, create them
        while len(paras) < 2:
            new_p = tf.add_paragraph()
            paras.append(new_p)

        # Replace first two meaningful paragraphs
        if ticket_bullet:
            _set_paragraph_text(paras[0], ticket_bullet)
            changes.append(f"  SUMM | Slide {slide_idx+1}: '{ticket_bullet}'")
        else:
            changes.append(f"  SUMM | Slide {slide_idx+1}: {label} ticket - no valid trend data (check missing/non-adjacent months)")

        if ageing_bullet:
            _set_paragraph_text(paras[1], ageing_bullet)
            changes.append(f"  SUMM | Slide {slide_idx+1}: '{ageing_bullet}'")
        else:
            changes.append(f"  SUMM | Slide {slide_idx+1}: {label} ageing - no valid trend data (check missing/non-adjacent months)")

        # Clear any extra old paragraphs beyond the first 2 meaningful ones
        for p in paras[2:]:
            _set_paragraph_text(p, "")

    return changes


# ==============================================================
# PROCESSING ENGINE v3 - HANDLES PICTURE, CHART, AND TABLE
# ==============================================================
def process_monthly_report(pdf_bytes, pptx_bytes, sel_month, sel_year):
    """
    Full automation engine.
    """
    from collections import defaultdict

    # -- Phase 1: Extract PDF page images --
    pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
    pdf_images = {}
    for pn in range(len(pdf)):
        page = pdf.load_page(pn)
        pix = page.get_pixmap(matrix=fitz.Matrix(4, 4))
        pdf_images[pn + 1] = pix.tobytes("png")
    log = [f"INFO | PDF loaded: {len(pdf)} pages at 300 DPI"]

    # -- Phase 2: Extract trend data (Smart Spatial Extraction) --
    sr_trend  = _normalize_chart_order(_extract_with_fallback(pdf, 3, "sr_trend", log))
    sr_ageing = _normalize_chart_order(_extract_with_fallback(pdf, 4, "sr_ageing", log))
    inc_trend = _normalize_chart_order(_extract_with_fallback(pdf, 10, "inc_trend", log))
    inc_ageing = _normalize_chart_order(_extract_with_fallback(pdf, 11, "inc_ageing", log))

    sr_table_data  = _extract_root_cause_table(pdf, 5,  "SR",  log)
    inc_table_data = _extract_root_cause_table(pdf, 12, "INC", log)

    # Debug: show final chronological order after normalization
    for key, vals in [("sr_trend", sr_trend), ("sr_ageing", sr_ageing),
                       ("inc_trend", inc_trend), ("inc_ageing", inc_ageing)]:
        if vals:
            log.append(f"DEBUG | FINAL ORDER ({key}): {[(m, int(v)) for m, v in vals]}")

    # Detailed logging for verification
    for lbl, vals in [
        ("SR Trend (p3)",    sr_trend),
        ("SR Ageing (p4)",   sr_ageing),
        ("INC Trend (p10)",  inc_trend),
        ("INC Ageing (p11)", inc_ageing),
    ]:
        if vals:
            pairs = ", ".join(f"{m}={int(v)}" for m, v in vals)
            log.append(f"INFO | Extract {lbl}: {pairs}")

    # -- Phase 3: Auto-crop module list pages --
    for pg in (7, 14):
        if pg in pdf_images:
            old = len(pdf_images[pg])
            pdf_images[pg] = _auto_crop_bottom(pdf_images[pg])
            nw = len(pdf_images[pg])
            if nw < old:
                log.append(f"INFO | PDF p{pg}: cropped ({old//1024}KB -> {nw//1024}KB)")

    # -- Phase 4: Open PPTX --
    prs = Presentation(io.BytesIO(pptx_bytes))
    log.append(f"INFO | PPTX: {len(prs.slides)} slides")

    LOGO_TOP = 0.5
    LOGO_H   = 0.6

    if len(prs.slides) > 9:
        slide = prs.slides[9]
        for sh in slide.shapes:
            try:
                t = _emu_to_inches(sh.top)
                h = _emu_to_inches(sh.height)
                is_logo = t < LOGO_TOP and h < LOGO_H
                log.append(f"DEBUG | Slide10 shape: type={sh.shape_type} name='{sh.name}' "
                           f"L={_emu_to_inches(sh.left):.2f} T={t:.2f} "
                           f"W={_emu_to_inches(sh.width):.2f} H={h:.2f} logo_filtered={is_logo}")
            except Exception as ex:
                log.append(f"DEBUG | Slide10 shape error: {ex}")

    # -- Phase 5: Replacement map --
    # (slide_idx_0, pdf_page, strategy, hint, label)
    MAP = [
        (2,  2, "BBOX", (0.314, 2.276, 3.782, 1.200), "SR SLA Performance"),
        (2,  3, "BBOX", (4.570, 2.276, 8.307, 3.286), "SR Ticket Trend"),
        (3,  4, "CHART", None, "SR Ageing Trend"),
        (5,  6, "BBOX", (0.488, 2.641, 6.533, 2.449), "SR Category Distribution"),
        (5,  7, "BBOX_LIST", (7.346, 1.448, 3.477, 5.459), "SR Module Ticket List"),
        (6,  8, "BBOX", (0.404, 2.406, 3.240, 1.313), "INC Response SLA"),
        (6,  9, "BBOX", (0.315, 4.294, 3.417, 1.365), "INC Resolution SLA"),
        (6, 10, "BBOX", (3.932, 2.406, 9.124, 3.567), "INC Ticket Trend"),
        (7, 11, "CHART", None, "INC Ageing Trend"),
        (9, 13, "BBOX", (0.860, 2.650, 6.090, 2.520), "INC Category Distribution"),
        (9, 14, "BBOX_LIST", (7.560, 2.650, 5.000, 2.520), "INC Module Ticket List"),
    ]

    replaced  = 0

    groups = defaultdict(list)
    for e in MAP:
        groups[e[0]].append(e)

    for si in sorted(groups.keys()):
        entries = groups[si]
        if si >= len(prs.slides):
            log.append(f"WARN | Slide {si+1} missing")
            continue

        slide = prs.slides[si]

        # Collect non-logo PICTURE shapes
        pics = []
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            if _emu_to_inches(sh.top) < LOGO_TOP and _emu_to_inches(sh.height) < LOGO_H:
                continue
            pics.append(sh)

        used = set()
        jobs = []  # (entry, left, top, width, height, shape_to_delete)

        for entry in entries:
            _, pp, strat, hint, label = entry
            if pp not in pdf_images:
                log.append(f"WARN | PDF p{pp} missing. Skip '{label}'.")
                continue

            target = None

            if strat in ("BBOX", "BBOX_LIST"):
                tl, tt, tw, th = hint
                best, bdist = None, float('inf')
                for sh in pics:
                    if id(sh) in used:
                        continue
                    d = (((_emu_to_inches(sh.left) - tl)**2) +
                         ((_emu_to_inches(sh.top)  - tt)**2) +
                         ((_emu_to_inches(sh.width) - tw)**2) +
                         ((_emu_to_inches(sh.height)- th)**2)) ** 0.5
                    if d < bdist:
                        bdist, best = d, sh
                if best and bdist < 3.0:
                    target = best
                    used.add(id(best))
                else:
                    log.append(f"WARN | Slide {si+1}: No PICTURE for '{label}' (dist={bdist:.2f})")

            elif strat == "CHART":
                for sh in slide.shapes:
                    try:
                        if sh.shape_type == MSO_SHAPE_TYPE.CHART:
                            target = sh
                            break
                    except:
                        pass
                if not target:
                    ba = 0
                    for sh in slide.shapes:
                        try:
                            st = sh.shape_type
                        except:
                            continue
                        if st == MSO_SHAPE_TYPE.PICTURE or st == 17:
                            continue
                        ti = _emu_to_inches(sh.top)
                        hi = _emu_to_inches(sh.height)
                        if ti < LOGO_TOP and hi < LOGO_H:
                            continue
                        if _emu_to_inches(sh.width) < 1.0 or hi < 0.5:
                            continue
                        a = sh.width * sh.height
                        if a > ba:
                            ba, target = a, sh
                if not target:
                    log.append(f"WARN | Slide {si+1}: No CHART for '{label}'")

            elif strat == "TABLE":
                # Find the native PPTX table shape
                for sh in slide.shapes:
                    try:
                        if sh.has_table:
                            target = sh
                            break
                    except:
                        pass
                if not target:
                    for sh in slide.shapes:
                        try:
                            if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                                target = sh
                                break
                        except:
                            pass
                if not target:
                    log.append(f"WARN | Slide{si+1}: No TABLE for '{label}'")

            elif strat == "POS":
                # Collect ALL non-logo shapes (not just PICTURE type),
                # since split placeholders may have a different shape_type
                all_non_logo = []
                for sh in slide.shapes:
                    try:
                        if _emu_to_inches(sh.top) < LOGO_TOP and _emu_to_inches(sh.height) < LOGO_H:
                            continue
                        # Skip text-only shapes
                        if hasattr(sh, 'text_frame') and not hasattr(sh, 'image'):
                            try:
                                _ = sh.shape_type
                                if sh.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX,):
                                    continue
                            except:
                                pass
                        all_non_logo.append(sh)
                    except:
                        continue
                
                avail = sorted([s for s in all_non_logo if id(s) not in used], key=lambda s: s.left)
                idx = hint
                if idx < len(avail):
                    target = avail[idx]
                    used.add(id(target))
                    log.append(f"INFO | Slide {si+1}: POS#{idx} -> shape type={target.shape_type} name='{target.name}'")
                else:
                    log.append(f"WARN | Slide {si+1}: Only {len(avail)} shapes, need #{idx} for '{label}'")

            if target:
                jobs.append((entry, target.left, target.top, target.width, target.height, target))

        # Delete old shapes first, then insert new images
        for _, _, _, _, _, sh in jobs:
            try:
                sh._element.getparent().remove(sh._element)
            except Exception as ex:
                log.append(f"WARN | Delete failed: {ex}")

        for entry, left, top, width, height, _ in jobs:
            pp, label, strat = entry[1], entry[4], entry[2]
            img_bytes = pdf_images[pp]

            # -- Smart placement based on strategy --
            if strat in ("CHART", "TABLE"):
                iw, ih = _get_image_size(img_bytes)
                ins_l, ins_t, ins_w, ins_h = _calc_fit_rect(iw, ih, left, top, width, height)
            elif strat == "BBOX_LIST":
                iw, ih = _get_image_size(img_bytes)
                ins_l, ins_t, ins_w, ins_h = _calc_list_rect(iw, ih, left, top, width)
            else:
                ins_l, ins_t, ins_w, ins_h = left, top, width, height

            slide.shapes.add_picture(io.BytesIO(img_bytes), ins_l, ins_t, ins_w, ins_h)
            b = (_emu_to_inches(ins_l), _emu_to_inches(ins_t),
                 _emu_to_inches(ins_w), _emu_to_inches(ins_h))
            log.append(f"  OK | Slide {entry[0]+1}: '{label}' [{strat}] <- PDF p{pp} "
                       f"(L={b[0]:.2f} T={b[1]:.2f} W={b[2]:.2f} H={b[3]:.2f})")
            replaced += 1

    # -- Phase 5b: Fill native PPTX tables with extracted data --
    log.append("INFO | Filling native PPTX tables...")

    # SR table on Slide 5 (index 4)
    if sr_table_data and len(prs.slides) > 4:
        slide5 = prs.slides[4]
        sr_tbl_shape = None
        for sh in slide5.shapes:
            try:
                if sh.has_table:
                    sr_tbl_shape = sh
                    break
            except:
                pass
        if sr_tbl_shape:
            _fill_pptx_table(sr_tbl_shape, sr_table_data, "SR", sel_month, sel_year, log)
            log.append("INFO | SR table filled successfully.")
        else:
            log.append("WARN | SR table shape not found on Slide 5.")
    else:
        log.append("WARN | SR table: no extracted data or slide missing.")

    # INC table on Slide 9 (index 8)
    if inc_table_data and len(prs.slides) > 8:
        slide9 = prs.slides[8]
        inc_tbl_shape = None
        for sh in slide9.shapes:
            try:
                if sh.has_table:
                    inc_tbl_shape = sh
                    break
            except:
                pass
        if inc_tbl_shape:
            _fill_pptx_table(inc_tbl_shape, inc_table_data, "INC", sel_month, sel_year, log)
            log.append("INFO | INC table filled successfully.")
        else:
            log.append("WARN | INC table shape not found on Slide 9.")
    else:
        log.append("WARN | INC table: no extracted data or slide missing.")

    # -- Phase 6: Date replacement --
    log.append("INFO | Updating dates...")
    dc = update_dates_in_pptx(prs, sel_month, str(sel_year))
    log.extend(dc)
    log.append(f"INFO | {len(dc)} date updates.")

    # -- Phase 7: Summary bullets (SAFE) --
    log.append("INFO | Updating summary on slides 4 and 8...")
    sc = update_summary_bullets(prs, sel_month, sel_year, sr_trend, sr_ageing, inc_trend, inc_ageing)
    log.extend(sc)

    # -- Phase 8: Save --
    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    log.append(f"INFO | Done: {replaced}/11 images replaced.")
    return out.read(), log, replaced, pdf_images


# -- Resolve PPTX template --
resolved_pptx = None
if uploaded_template is not None:
    resolved_pptx = "uploaded"
elif pptx_available_locally:
    resolved_pptx = "local"


# ==============================================================
# GENERATE SECTION
# ==============================================================
if pdf_file and resolved_pptx:
    st.markdown('<p class="section-label">Step 2 - Validation and Generation</p>', unsafe_allow_html=True)

    pdf_file.seek(0)
    try:
        preview_pdf = fitz.open(stream=pdf_file.read(), filetype="pdf")
        pdf_page_count = len(preview_pdf)
        pdf_file.seek(0)
    except:
        pdf_page_count = "?"

    m1, m2, m3 = st.columns(3)
    with m1:
        st.metric("PDF Pages Detected", pdf_page_count)
    with m2:
        st.metric("Target Slides", "8 (Slides 3-10)")
    with m3:
        st.metric("Image Swaps", "11")

    with st.expander("PDF to PPT Mapping", expanded=False):
        st.markdown("""
| PDF | Content | Slide | Method | Notes |
|:---:|---|:---:|---|---|
| 1 | Landing | *Skip* | - | - |
| **2** | SR SLA | **3** | BBOX match | Left picture |
| **3** | SR Trend | **3** | BBOX match | Right picture |
| **4** | SR Ageing | **4** | Delete chart | Replaces native histogram |
| **5** | SR Root Cause | **5** | Native Table | Populates cells natively |
| **6** | SR Category % | **6** | BBOX match | Left picture |
| **7** | SR Module List | **6** | BBOX match | Right picture (auto-cropped) |
| **8** | INC Response SLA | **7** | BBOX match | Top-left picture |
| **9** | INC Resolution SLA | **7** | BBOX match | Bottom-left picture |
| **10** | INC Trend | **7** | BBOX match | Right picture |
| **11** | INC Ageing | **8** | Delete chart | Replaces native histogram |
| **12** | INC Root Cause | **9** | Native Table | Populates cells natively |
| **13** | INC Category % | **10** | BBOX match | Left picture |
| **14** | INC Module List | **10** | BBOX match | Right picture (auto-cropped) |
        """)
        st.info(
            "**Automated features:**\n"
            "- Module lists (pages 7 and 14) are auto-cropped to remove whitespace\n"
            "- Root Cause tables (pages 5 and 12) are mapped natively into PowerPoint tables\n"
            "- Dates are updated across all slides\n"
            "- Summary bullets on Slides 4 and 8 are auto-computed from extracted chart data\n"
            "- Corporate logo is protected and title text is never modified"
        )

    if pdf_page_count != "?" and pdf_page_count < 14:
        st.warning(f"Expected 14 pages but found {pdf_page_count}. Some images may be missing.")

    st.markdown("")

    if st.button("Generate Monthly Report", use_container_width=True, type="primary"):
        with st.spinner("Extracting visuals, replacing charts and tables, updating text..."):
            try:
                if uploaded_template is not None:
                    template_bytes = uploaded_template.read()
                else:
                    with open(TEMPLATE_PATH, "rb") as f:
                        template_bytes = f.read()

                out_bytes, build_logs, img_count, pdf_imgs = process_monthly_report(
                    pdf_file.read(), template_bytes, sel_month, int(sel_year)
                )

                if img_count == 11:
                    st.success(f"Complete - all {img_count}/11 images replaced and text updated.")
                elif img_count > 0:
                    st.warning(f"Partial - {img_count}/11 images replaced. Review the build log.")
                else:
                    st.error("No images were replaced. Please check your input files.")

                with st.expander("Build Log", expanded=(img_count < 11)):
                    for msg in build_logs:
                        if "WARN" in msg:
                            st.warning(msg)
                        elif msg.startswith("  OK") or msg.startswith("  DATE") or msg.startswith("  SUMM") or msg.startswith("  TABLE"):
                            st.text(msg)
                        else:
                            st.info(msg)

                # -- Preview --
                if img_count > 0 and pdf_imgs:
                    st.markdown("")
                    st.markdown('<p class="section-label">Preview - Images Placed Per Slide</p>', unsafe_allow_html=True)
                    PREVIEW = [
                        ("Slide 3 - SR SLA and Trend",               [2, 3]),
                        ("Slide 4 - SR Ageing (chart replaced)",      [4]),
                        ("Slide 5 - SR Root Cause (table updated)",   [5]),
                        ("Slide 6 - SR Category and Module List",     [6, 7]),
                        ("Slide 7 - INC SLA and Trend",              [8, 9, 10]),
                        ("Slide 8 - INC Ageing (chart replaced)",     [11]),
                        ("Slide 9 - INC Root Cause (table updated)",  [12]),
                        ("Slide 10 - INC Category and Module List",   [13, 14]),
                    ]
                    for title, pages in PREVIEW:
                        avail = [p for p in pages if p in pdf_imgs]
                        if not avail:
                            continue
                        with st.expander(title, expanded=False):
                            cols = st.columns(len(avail))
                            for col, pg in zip(cols, avail):
                                with col:
                                    st.caption(f"PDF Page {pg}")
                                    st.image(pdf_imgs[pg], use_container_width=True)

                # -- Download --
                if img_count > 0:
                    st.markdown("")
                    st.markdown('<p class="section-label">Step 3 - Download</p>', unsafe_allow_html=True)
                    st.info(
                        "**What was automatically updated:**\n"
                        "- All 11 dashboard images replaced\n"
                        "- Charts on Slides 4 and 8 replaced with PDF images\n"
                        "- Tables on Slides 5 and 9 populated natively with PDF data\n"
                        "- Dates updated across all slides\n"
                        "- Summary bullets on Slides 4 and 8 computed from chart data (if extractable)\n\n"
                        "Review summary text on Slides 4 and 8 if the build log shows warnings."
                    )
                    st.download_button(
                        label="Download Final Report (.pptx)",
                        data=out_bytes,
                        file_name=f"Monthly_Report_{sel_month}_{sel_year}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )

            except Exception as e:
                st.error(f"An error occurred: {e}")
                st.code(traceback.format_exc(), language="text")

elif not resolved_pptx and pdf_file:
    st.warning("No PPTX template found. Please upload one in the sidebar.")

else:
    st.markdown("""
    <div style="
        text-align: center; padding: 60px 40px;
        background: linear-gradient(180deg, #FFFFFF 0%, #F0FAFA 100%);
        border: 1px solid #E2E8F0; border-top: 4px solid #00B1A9;
        border-radius: 16px; margin-top: 8px;
    ">
        <h2 style="color: #1A202C; font-weight: 800; margin: 0 0 10px 0;">Upload Required</h2>
        <p style="color: #718096; max-width: 520px; margin: 0 auto; line-height: 1.7; font-size: 0.95rem;">
            Select the Report Month and Year, then upload your Power BI PDF export in the sidebar.
        </p>
    </div>
    """, unsafe_allow_html=True)